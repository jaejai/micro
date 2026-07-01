"""PowerPoint report builder — ports notebook section §13.

Captures figures (from plotting.py) in memory and assembles a styled deck.
REPORT mode: "all" | "ebsd" | "odf".
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

from .config import Config
from .microstructure import MicroResult
from .odf import ODFResult
from . import plotting as P

# palette
DEEP = RGBColor(0x06, 0x5A, 0x82); TEAL = RGBColor(0x1C, 0x72, 0x93); MID = RGBColor(0x21, 0x29, 0x5C)
ICE = RGBColor(0xCA, 0xDC, 0xFC); LIGHT = RGBColor(0xEC, 0xF2, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); MUTE = RGBColor(0x5B, 0x6B, 0x7A)
HF = "Cambria"; BF = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)   # overridden per-build from cfg


def _grab(fig, dpi=150):
    # figures use constrained_layout; don't combine with bbox_inches='tight'
    buf = io.BytesIO(); fig.savefig(buf, format="png", dpi=dpi); buf.seek(0); return buf


def _solid(sh, c):
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background(); sh.shadow.inherit = False


def _tbx(s, l, t, w, h, a=MSO_ANCHOR.TOP):
    f = s.shapes.add_textbox(l, t, w, h).text_frame; f.word_wrap = True; f.vertical_anchor = a
    f.margin_left = f.margin_right = Inches(0.05); f.margin_top = f.margin_bottom = Inches(0.02); return f


def _par(p, t, sz, c, b=False, fn=BF, al=PP_ALIGN.LEFT, it=False):
    p.text = t; p.alignment = al; r = p.runs[0]
    r.font.size = Pt(sz); r.font.bold = b; r.font.italic = it; r.font.color.rgb = c; r.font.name = fn


def _fit(s, buf, bl, bt, bw, bh):
    buf.seek(0); w, h = Image.open(buf).size; buf.seek(0); ar = w / h
    if ar > bw / bh: nw = bw; nh = int(bw / ar)
    else: nh = bh; nw = int(bh * ar)
    s.shapes.add_picture(buf, bl + (bw - nw) // 2, bt + (bh - nh) // 2, nw, nh)


def _title_slide(prs, title, subtitle, metrics):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _solid(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH), MID)
    _solid(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SH), TEAL)
    _par(_tbx(s, Inches(0.85), Inches(0.85), Inches(11.6), Inches(1.0)).paragraphs[0], title, 40, WHITE, b=True, fn=HF)
    _par(_tbx(s, Inches(0.88), Inches(1.95), Inches(11.6), Inches(0.5)).paragraphs[0], subtitle, 18, ICE, it=True)
    x0, y0, cw, ch = Inches(0.85), Inches(2.95), Inches(3.75), Inches(1.55)
    for i, (v, l) in enumerate(metrics[:6]):
        r, c = divmod(i, 3); x = x0 + c * (cw + Inches(0.2)); y = y0 + r * (ch + Inches(0.2))
        _solid(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch), DEEP)
        _par(_tbx(s, x, y + Inches(0.18), cw, Inches(0.75), MSO_ANCHOR.MIDDLE).paragraphs[0], v, 28, WHITE, b=True, fn=HF, al=PP_ALIGN.CENTER)
        _par(_tbx(s, x, y + Inches(0.95), cw, Inches(0.45)).paragraphs[0], l.upper(), 10.5, ICE, al=PP_ALIGN.CENTER)
    return s


def _heading(s, text, size=30):
    _solid(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH), WHITE)
    _solid(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH), TEAL)
    _par(_tbx(s, Inches(0.55), Inches(0.35), Inches(12.4), Inches(0.8)).paragraphs[0], text, size, DEEP, b=True, fn=HF)


def add_ebsd_slides(prs, cfg: Config, res: MicroResult):
    dpi = cfg.fig_dpi; nb = cfg.hist_bins
    figs = {
        "iq": _grab(P.fig_iq(res), dpi), "ci": _grab(P.fig_ci(res), dpi),
        "ipf": _grab(P.fig_ipf(cfg, res), dpi), "grains": _grab(P.fig_grains(cfg, res), dpi),
        "gb": _grab(P.fig_gb(cfg, res), dpi), "ipfgb": _grab(P.fig_ipf_hagb(cfg, res), dpi),
        "gs_count": _grab(P.fig_grain_size_count(res, nb), dpi),
        "gs_frac": _grab(P.fig_grain_size_frac(res, nb), dpi),
    }
    _title_slide(prs, "EBSD Microstructure Analysis",
                 f"DP590 dual-phase steel  -  {cfg.crystal_sym}  -  {res.nx}x{res.ny} @ {res.step:.3f} um/px",
                 [(f"{res.n_grains}", f"Grains (>= {cfg.min_grain_px} px)"), (f"{res.G_e2627:.1f}", "ASTM E2627 G"),
                  (f"{res.d_num:.2f} um", "Number-avg. dia."), (f"{res.d_w:.2f} um", "Area-weighted dia."),
                  (f"{res.g_area_um2.mean():.2f} um2", "Mean grain area"), (f"{cfg.hagb_angle:.0f}deg", "HAGB threshold")])
    # maps grid
    s = prs.slides.add_slide(prs.slide_layouts[6]); _heading(s, "Orientation & Quality Maps")
    cw = Inches(3.0); x = Inches(0.6); pt = Inches(1.4); ph = Inches(5.2)
    for key, lab in [("iq", "IQ"), ("ci", "CI"), ("ipf", f"IPF [{cfg.dir_str}]"), ("grains", f"Grains ({res.n_grains})")]:
        _solid(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, pt, cw, ph), LIGHT)
        _fit(s, figs[key], x + Inches(0.1), pt + Inches(0.1), cw - Inches(0.2), ph - Inches(0.55))
        _par(_tbx(s, x, pt + ph - Inches(0.5), cw, Inches(0.4), MSO_ANCHOR.MIDDLE).paragraphs[0], lab, 14, DEEP, b=True, al=PP_ALIGN.CENTER)
        x += cw + Inches(0.18)
    # boundaries
    s = prs.slides.add_slide(prs.slide_layouts[6]); _heading(s, "Grain Boundaries")
    for key, lab, x in [("gb", f"GB map (LAGB {cfg.lagb_angle:.0f}deg/HAGB {cfg.hagb_angle:.0f}deg)", Inches(1.4)),
                        ("ipfgb", f"IPF [{cfg.dir_str}] + HAGB", Inches(7.1))]:
        _solid(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(4.8), Inches(5.2)), LIGHT)
        _fit(s, figs[key], x + Inches(0.1), Inches(1.5), Inches(4.6), Inches(4.55))
        _par(_tbx(s, x, Inches(6.1), Inches(4.8), Inches(0.4), MSO_ANCHOR.MIDDLE).paragraphs[0], lab, 14, DEEP, b=True, al=PP_ALIGN.CENTER)
    # grain size
    s = prs.slides.add_slide(prs.slide_layouts[6]); _heading(s, "Grain Size - Count & Area Fraction", 28)
    _fit(s, figs["gs_count"], Inches(0.55), Inches(1.05), Inches(12.2), Inches(2.95))
    _fit(s, figs["gs_frac"], Inches(0.55), Inches(4.15), Inches(12.2), Inches(2.95))


def add_odf_slides(prs, cfg: Config, res: MicroResult, odf: ODFResult):
    dpi = cfg.fig_dpi
    figs = {"sections": _grab(P.fig_odf_sections(cfg, odf), dpi), "fibers": _grab(P.fig_fibers(odf), dpi)}
    _title_slide(prs, "Crystallographic Texture (ODF)",
                 f"DP590 {cfg.lattice} ferrite  -  GSH L_max={cfg.l_max}  -  {len(odf.eulers_odf):,} orientations",
                 [(f"{odf.J:.2f}", "Texture index J"), (f"{odf.odf.max():.2f}", "ODF max [mrd]"),
                  (f"{odf.f_alpha.max():.2f}", "alpha-fiber max [mrd]"), (f"{odf.f_gamma.max():.2f}", "gamma-fiber max [mrd]"),
                  (f"{len(odf.n_states)}", "GSH modes"), (f"L={cfg.l_max}", "Bandwidth")])
    s = prs.slides.add_slide(prs.slide_layouts[6]); _heading(s, "ODF - phi2 Sections")
    _fit(s, figs["sections"], Inches(0.55), Inches(1.15), Inches(12.2), Inches(5.4))
    _par(_tbx(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.5)).paragraphs[0],
         f"Max density at {odf.odf_max_loc}  -  white markers: ideal {cfg.lattice} components", 12, MUTE, al=PP_ALIGN.CENTER)
    s = prs.slides.add_slide(prs.slide_layouts[6]); _heading(s, "Fiber Texture (alpha & gamma)")
    _fit(s, figs["fibers"], Inches(0.55), Inches(1.25), Inches(12.2), Inches(3.8))
    import numpy as np
    for v, l, x in [(f"{odf.f_alpha.max():.2f} mrd", f"alpha-fiber peak (Phi={odf.Phi_line[odf.f_alpha.argmax()]:.0f}deg)", Inches(2.3)),
                    (f"{odf.f_gamma.max():.2f} mrd", f"gamma-fiber peak (phi1={odf.phi1_line[odf.f_gamma.argmax()]:.0f}deg)", Inches(7.3))]:
        _solid(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.5), Inches(3.7), Inches(1.3)), DEEP)
        _par(_tbx(s, x, Inches(5.62), Inches(3.7), Inches(0.75), MSO_ANCHOR.MIDDLE).paragraphs[0], v, 26, WHITE, b=True, fn=HF, al=PP_ALIGN.CENTER)
        _par(_tbx(s, x, Inches(6.4), Inches(3.7), Inches(0.4)).paragraphs[0], l.upper(), 10.5, ICE, al=PP_ALIGN.CENTER)


def build_report(cfg: Config, res: MicroResult, odf: ODFResult, outfile: str, mode: str = "all", log=print) -> str:
    global SW, SH
    SW, SH = Inches(cfg.slide_w_in), Inches(cfg.slide_h_in)
    prs = Presentation(); prs.slide_width, prs.slide_height = SW, SH
    if mode in ("all", "ebsd"):
        add_ebsd_slides(prs, cfg, res)
    if mode in ("all", "odf"):
        if odf is None:
            raise ValueError("ODF result required for mode '%s'" % mode)
        add_odf_slides(prs, cfg, res, odf)
    prs.save(outfile)
    log(f"Saved {outfile}  ({len(prs.slides._sldIdLst)} slides, mode='{mode}')")
    return outfile
